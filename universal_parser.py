"""
범용 문서 파서 - 다양한 파일 형식 지원
세계 최고 수준의 문서 변환 엔진
"""

import os
import logging
import mimetypes
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from abc import ABC, abstractmethod

logger = logging.getLogger(__name__)


class DocumentParser(ABC):
    """문서 파서 기본 클래스"""

    @abstractmethod
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        """이 파서가 해당 파일을 처리할 수 있는지 확인"""
        pass

    @abstractmethod
    def parse(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        """파일 파싱 및 데이터 추출"""
        pass

    @abstractmethod
    def get_supported_formats(self) -> List[str]:
        """지원하는 파일 형식 목록"""
        pass


class PDFParser(DocumentParser):
    """PDF 파서 (기존 PDFConverter 활용)"""

    def __init__(self):
        from pdf_converter import PDFConverter
        self.converter = PDFConverter()

    def can_parse(self, file_path: str, mime_type: str) -> bool:
        return mime_type == 'application/pdf' or file_path.lower().endswith('.pdf')

    def parse(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        content_type = options.get('content_type', 'general')
        result = self.converter.extract_from_pdf(file_path, content_type=content_type)
        if result:
            result['format'] = 'pdf'
        return result or {}

    def get_supported_formats(self) -> List[str]:
        return ['pdf', 'application/pdf']


class ImageParser(DocumentParser):
    """이미지 파서 (JPG, PNG, TIFF 등)"""

    def __init__(self):
        try:
            from vision_ocr import VisionOCR
            self.ocr = VisionOCR()
        except:
            self.ocr = None
            logger.warning("Vision OCR 초기화 실패 - 이미지 파싱 불가")

    def can_parse(self, file_path: str, mime_type: str) -> bool:
        image_types = ['image/jpeg', 'image/jpg', 'image/png', 'image/tiff', 'image/bmp', 'image/webp']
        extensions = ['.jpg', '.jpeg', '.png', '.tiff', '.tif', '.bmp', '.webp']
        return mime_type in image_types or any(file_path.lower().endswith(ext) for ext in extensions)

    def parse(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        if not self.ocr or not self.ocr.client:
            return {"error": "OCR 사용 불가"}

        try:
            import base64
            from PIL import Image
            import io

            # 이미지 로드 및 리사이징
            img = Image.open(file_path)
            if img.mode == 'RGBA':
                img = img.convert('RGB')

            # 최적 크기로 리사이징
            max_size = 2000
            if max(img.width, img.height) > max_size:
                ratio = max_size / max(img.width, img.height)
                new_size = (int(img.width * ratio), int(img.height * ratio))
                img = img.resize(new_size, Image.LANCZOS)

            # Base64 인코딩
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=90)
            base64_img = base64.b64encode(buffer.getvalue()).decode()

            # OCR 실행
            content_type = options.get('content_type', 'general')
            if content_type == 'election':
                ocr_result = self.ocr.extract_election_info(base64_img)
                text = ocr_result.get('text', '')
                structured = ocr_result.get('structured', {})
            else:
                text = self.ocr.extract_text_from_image(base64_img)
                structured = {}

            return {
                'format': 'image',
                'filename': Path(file_path).name,
                'page_count': 1,
                'pages': [{
                    'page_number': 1,
                    'text': text,
                    'image': f'data:image/jpeg;base64,{base64_img}',
                    'width': img.width,
                    'height': img.height
                }],
                'structured_data': structured,
                'is_image_based': True,
                'ocr_used': True
            }

        except Exception as e:
            logger.error(f"이미지 파싱 실패: {str(e)}", exc_info=True)
            return {"error": str(e)}

    def get_supported_formats(self) -> List[str]:
        return ['jpg', 'jpeg', 'png', 'tiff', 'tif', 'bmp', 'webp']


class WordParser(DocumentParser):
    """Microsoft Word 문서 파서 (.docx, .doc)"""

    def can_parse(self, file_path: str, mime_type: str) -> bool:
        word_types = ['application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                     'application/msword']
        extensions = ['.docx', '.doc']
        return mime_type in word_types or any(file_path.lower().endswith(ext) for ext in extensions)

    def parse(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        try:
            try:
                import docx
            except ImportError:
                return {"error": "python-docx 패키지 필요: pip install python-docx"}

            doc = docx.Document(file_path)

            # 텍스트 추출
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # 표 추출
            tables_data = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables_data.append(table_data)

            full_text = '\n'.join(paragraphs)

            return {
                'format': 'word',
                'filename': Path(file_path).name,
                'page_count': len(doc.sections),
                'pages': [{
                    'page_number': 1,
                    'text': full_text,
                    'paragraphs': paragraphs,
                    'tables': tables_data
                }],
                'is_image_based': False,
                'ocr_used': False
            }

        except Exception as e:
            logger.error(f"Word 문서 파싱 실패: {str(e)}", exc_info=True)
            return {"error": str(e)}

    def get_supported_formats(self) -> List[str]:
        return ['docx', 'doc']


class PowerPointParser(DocumentParser):
    """PowerPoint 프레젠테이션 파서 (.pptx, .ppt)"""

    def can_parse(self, file_path: str, mime_type: str) -> bool:
        ppt_types = ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    'application/vnd.ms-powerpoint']
        extensions = ['.pptx', '.ppt']
        return mime_type in ppt_types or any(file_path.lower().endswith(ext) for ext in extensions)

    def parse(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        try:
            try:
                from pptx import Presentation
            except ImportError:
                return {"error": "python-pptx 패키지 필요: pip install python-pptx"}

            prs = Presentation(file_path)

            slides_data = []
            for idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_text.append(shape.text)

                slides_data.append({
                    'page_number': idx,
                    'text': '\n'.join(slide_text),
                    'shape_count': len(slide.shapes)
                })

            all_text = '\n\n'.join([s['text'] for s in slides_data])

            return {
                'format': 'powerpoint',
                'filename': Path(file_path).name,
                'page_count': len(prs.slides),
                'pages': slides_data,
                'metadata': {
                    'slide_count': len(prs.slides),
                },
                'is_image_based': False,
                'ocr_used': False
            }

        except Exception as e:
            logger.error(f"PowerPoint 파싱 실패: {str(e)}", exc_info=True)
            return {"error": str(e)}

    def get_supported_formats(self) -> List[str]:
        return ['pptx', 'ppt']


class ExcelParser(DocumentParser):
    """Excel 스프레드시트 파서 (.xlsx, .xls)"""

    def can_parse(self, file_path: str, mime_type: str) -> bool:
        excel_types = ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                      'application/vnd.ms-excel']
        extensions = ['.xlsx', '.xls']
        return mime_type in excel_types or any(file_path.lower().endswith(ext) for ext in extensions)

    def parse(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        try:
            try:
                import openpyxl
            except ImportError:
                return {"error": "openpyxl 패키지 필요: pip install openpyxl"}

            wb = openpyxl.load_workbook(file_path, data_only=True)

            sheets_data = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]

                # 데이터 추출
                data = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        data.append(list(row))

                # 텍스트 변환
                text_lines = []
                for row in data:
                    text_lines.append('\t'.join([str(cell) if cell is not None else '' for cell in row]))

                sheets_data.append({
                    'sheet_name': sheet_name,
                    'data': data,
                    'text': '\n'.join(text_lines),
                    'rows': len(data),
                    'columns': max(len(row) for row in data) if data else 0
                })

            return {
                'format': 'excel',
                'filename': Path(file_path).name,
                'page_count': len(wb.sheetnames),
                'pages': sheets_data,
                'metadata': {
                    'sheet_count': len(wb.sheetnames),
                    'sheet_names': wb.sheetnames
                },
                'is_image_based': False,
                'ocr_used': False
            }

        except Exception as e:
            logger.error(f"Excel 파싱 실패: {str(e)}", exc_info=True)
            return {"error": str(e)}

    def get_supported_formats(self) -> List[str]:
        return ['xlsx', 'xls', 'csv']


class HWPParser(DocumentParser):
    """한글(HWP) 파일 파서 (.hwp, .hwpx)"""

    def can_parse(self, file_path: str, mime_type: str) -> bool:
        hwp_types = ['application/x-hwp', 'application/haansofthwp']
        extensions = ['.hwp', '.hwpx']
        return mime_type in hwp_types or any(file_path.lower().endswith(ext) for ext in extensions)

    def parse(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        """HWP 파일 파싱"""
        try:
            # HWPX (한글 2014+) - XML 기반
            if file_path.lower().endswith('.hwpx'):
                return self._parse_hwpx(file_path, options)
            # HWP (한글 97-2010) - OLE 기반
            else:
                return self._parse_hwp(file_path, options)

        except Exception as e:
            logger.error(f"HWP 파싱 실패: {str(e)}", exc_info=True)
            # HWP 파싱 실패 시 OCR로 폴백
            return self._parse_hwp_with_ocr(file_path, options)

    def _parse_hwp(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        """HWP 파일 파싱 (olefile 사용)"""
        try:
            import olefile
            import struct

            # HWP는 OLE 복합 문서 형식
            ole = olefile.OleFileIO(file_path)

            text_content = []

            # HWP 파일 구조에서 텍스트 스트림 찾기
            # BodyText/Section* 스트림에 텍스트가 저장됨
            for entry in ole.listdir():
                entry_path = '/'.join(entry)

                # BodyText 섹션 찾기
                if 'BodyText' in entry_path and 'Section' in entry_path:
                    try:
                        stream = ole.openstream(entry)
                        data = stream.read()

                        # 간단한 텍스트 추출 (실제로는 복잡한 파싱 필요)
                        # UTF-16LE로 디코딩 시도
                        try:
                            text = data.decode('utf-16le', errors='ignore')
                            # 제어 문자 제거
                            text = ''.join(c for c in text if c.isprintable() or c in '\n\r\t')
                            if text.strip():
                                text_content.append(text.strip())
                        except:
                            pass
                    except Exception as e:
                        logger.debug(f"HWP 스트림 읽기 실패 ({entry_path}): {str(e)}")

            ole.close()

            all_text = '\n\n'.join(text_content)

            if not all_text.strip():
                # 텍스트 추출 실패 시 메시지 반환
                return {
                    "error": "HWP 파일 텍스트 추출 실패",
                    "suggestion": "HWPX 형식으로 저장하거나 PDF로 변환해주세요.",
                    "note": "HWP 97-2010 형식은 복잡한 구조로 인해 완벽한 파싱이 어렵습니다."
                }

            return {
                'format': 'hwp',
                'detected_format': 'hwp',
                'parser_used': 'olefile (기본)',
                'filename': Path(file_path).name,
                'page_count': 1,
                'pages': [{
                    'page_number': 1,
                    'text': all_text
                }],
                'metadata': {
                    'note': 'HWP 파일 기본 파싱. 완벽한 텍스트 추출을 위해 HWPX 또는 PDF 변환을 권장합니다.'
                },
                'is_image_based': False,
                'ocr_used': False
            }

        except ImportError:
            return {"error": "olefile 패키지 필요: pip install olefile"}
        except Exception as e:
            logger.error(f"HWP 파싱 실패: {str(e)}")
            return {
                "error": f"HWP 파싱 실패: {str(e)}",
                "suggestion": "HWPX 형식으로 저장하거나 PDF로 변환해주세요."
            }

    def _parse_hwpx(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        """HWPX 파일 파싱 (ZIP + XML 구조)"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET

            text_content = []

            with zipfile.ZipFile(file_path, 'r') as hwpx:
                # HWPX는 ZIP 파일 내부에 XML로 구성
                # Contents/section0.xml, section1.xml 등에 본문 텍스트
                for name in hwpx.namelist():
                    if 'section' in name.lower() and name.endswith('.xml'):
                        try:
                            xml_content = hwpx.read(name)
                            root = ET.fromstring(xml_content)

                            # 텍스트 노드 추출
                            for text_elem in root.iter():
                                if text_elem.text and text_elem.text.strip():
                                    text_content.append(text_elem.text.strip())
                        except Exception as e:
                            logger.debug(f"HWPX 섹션 파싱 실패 ({name}): {str(e)}")

            all_text = '\n'.join(text_content)

            if not all_text.strip():
                # 텍스트 추출 실패 시 OCR 시도
                return self._parse_hwp_with_ocr(file_path, options)

            return {
                'format': 'hwpx',
                'detected_format': 'hwpx',
                'parser_used': 'zipfile+xml',
                'filename': Path(file_path).name,
                'page_count': 1,
                'pages': [{
                    'page_number': 1,
                    'text': all_text
                }],
                'metadata': {
                    'file_type': 'hwpx',
                    'korean_version': '2014+'
                },
                'is_image_based': False,
                'ocr_used': False
            }

        except Exception as e:
            logger.error(f"HWPX 파싱 실패: {str(e)}")
            return self._parse_hwp_with_ocr(file_path, options)

    def _parse_hwp_with_ocr(self, file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
        """HWP 파일을 OCR로 처리 (폴백 방법)"""
        logger.info(f"HWP 파일 OCR 처리 시도: {Path(file_path).name}")

        # TODO: HWP → PDF 또는 HWP → 이미지 변환 후 OCR
        # 현재는 LibreOffice나 한컴오피스 API가 필요함

        return {
            "error": "HWP 파일 텍스트 추출 실패. hwp5 패키지를 설치하거나 PDF로 변환해주세요.",
            "fallback_required": True,
            "suggestion": "HWP 파일을 PDF로 변환한 후 다시 업로드해주세요."
        }

    def get_supported_formats(self) -> List[str]:
        return ['hwp', 'hwpx']


class UniversalDocumentParser:
    """범용 문서 파서 - 모든 형식 통합 처리"""

    def __init__(self):
        self.parsers: List[DocumentParser] = [
            PDFParser(),
            ImageParser(),
            WordParser(),
            PowerPointParser(),
            ExcelParser(),
            HWPParser(),  # 한글 파서 추가
        ]

        logger.info(f"범용 문서 파서 초기화 완료 - {len(self.parsers)}개 파서 등록 (HWP 지원 포함)")

    def get_supported_formats(self) -> Dict[str, List[str]]:
        """지원하는 모든 파일 형식 반환"""
        formats = {}
        for parser in self.parsers:
            parser_name = parser.__class__.__name__.replace('Parser', '')
            formats[parser_name] = parser.get_supported_formats()
        return formats

    def detect_file_type(self, file_path: str) -> Tuple[str, str]:
        """파일 타입 자동 감지"""
        mime_type, _ = mimetypes.guess_type(file_path)
        extension = Path(file_path).suffix.lower().lstrip('.')

        return mime_type or 'application/octet-stream', extension

    def parse_document(self, file_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        범용 문서 파싱

        Args:
            file_path: 파일 경로
            options: 파싱 옵션
                - content_type: 콘텐츠 타입 (election, lecture, church, general)
                - output_format: 출력 형식 (html, json, markdown 등)
                - custom_template: 사용자 정의 템플릿
                - language: 언어 (ko, en, ja, zh 등)

        Returns:
            파싱된 문서 데이터
        """
        if options is None:
            options = {}

        try:
            # 파일 존재 확인
            if not os.path.exists(file_path):
                return {"error": f"파일을 찾을 수 없습니다: {file_path}"}

            # 파일 타입 감지
            mime_type, extension = self.detect_file_type(file_path)
            logger.info(f"파일 타입 감지: {extension} ({mime_type})")

            # 적절한 파서 찾기
            for parser in self.parsers:
                if parser.can_parse(file_path, mime_type):
                    logger.info(f"{parser.__class__.__name__} 사용하여 파싱 시작")
                    result = parser.parse(file_path, options)

                    if result and 'error' not in result:
                        result['detected_format'] = extension
                        result['mime_type'] = mime_type
                        result['parser_used'] = parser.__class__.__name__
                        return result
                    else:
                        logger.warning(f"{parser.__class__.__name__} 파싱 실패")

            # 지원하지 않는 형식
            return {
                "error": f"지원하지 않는 파일 형식: {extension}",
                "detected_format": extension,
                "mime_type": mime_type,
                "supported_formats": self.get_supported_formats()
            }

        except Exception as e:
            logger.error(f"문서 파싱 중 오류: {str(e)}", exc_info=True)
            return {"error": str(e)}


# 싱글톤 인스턴스
_universal_parser = None

def get_universal_parser() -> UniversalDocumentParser:
    """범용 파서 싱글톤 인스턴스 가져오기"""
    global _universal_parser
    if _universal_parser is None:
        _universal_parser = UniversalDocumentParser()
    return _universal_parser
