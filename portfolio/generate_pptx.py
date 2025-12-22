"""
HTML 포트폴리오를 PPTX로 변환하는 스크립트
실행: python generate_pptx.py
출력: portfolio_2025.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 색상 정의
COLORS = {
    'primary': RGBColor(37, 99, 235),      # #2563EB
    'primary_dark': RGBColor(29, 78, 216), # #1D4ED8
    'secondary': RGBColor(124, 58, 237),   # #7C3AED
    'accent': RGBColor(16, 185, 129),      # #10B981
    'warning': RGBColor(245, 158, 11),     # #F59E0B
    'orange': RGBColor(234, 88, 12),       # #EA580C
    'dark': RGBColor(31, 41, 55),          # #1F2937
    'gray': RGBColor(107, 114, 128),       # #6B7280
    'light': RGBColor(243, 244, 246),      # #F3F4F6
    'white': RGBColor(255, 255, 255),
    'light_blue': RGBColor(100, 150, 255),
    'purple': RGBColor(167, 139, 250),
}

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = 'Malgun Gothic'
    p.alignment = align
    if color:
        p.font.color.rgb = color
    return txBox

def add_shape_with_text(slide, left, top, width, height, text, fill_color, text_color=None, font_size=14, bold=False):
    """색상 배경이 있는 도형에 텍스트 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = 'Malgun Gothic'
    p.alignment = PP_ALIGN.CENTER
    if text_color:
        p.font.color.rgb = text_color
    return shape

def create_slide_1_cover(prs):
    """슬라이드 1: 표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['primary_dark']
    background.line.fill.background()

    # 배지
    add_shape_with_text(slide, 3.5, 0.8, 3, 0.5, "2025 Digital Solutions Portfolio",
                        COLORS['white'], COLORS['primary_dark'], 11)

    # 메인 타이틀
    add_text_box(slide, 0.5, 1.8, 9, 1, "AI 기반 문서 변환 &\n통합 유통관리 솔루션",
                 36, True, COLORS['white'], PP_ALIGN.CENTER)

    # 서브타이틀
    add_text_box(slide, 0.5, 3.2, 9, 0.5, "PDF 자동 변환 · 다국어 학습 · B2B2C 커머스",
                 18, False, COLORS['white'], PP_ALIGN.CENTER)

    # 3개 솔루션 박스
    add_shape_with_text(slide, 0.8, 4.0, 2.8, 0.6, "📱 PDF 모바일 최적화",
                        COLORS['light_blue'], COLORS['white'], 12)
    add_shape_with_text(slide, 3.6, 4.0, 2.8, 0.6, "🌍 특수직군 외국어 학습",
                        COLORS['light_blue'], COLORS['white'], 12)
    add_shape_with_text(slide, 6.4, 4.0, 2.8, 0.6, "🛒 통합 유통관리",
                        COLORS['light_blue'], COLORS['white'], 12)

    # 연락처
    add_text_box(slide, 0.5, 5.0, 9, 0.4, "jmyangkr@gmail.com | 010-8665-8150",
                 14, False, COLORS['white'], PP_ALIGN.CENTER)

def create_slide_2_toc(prs):
    """슬라이드 2: 목차"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    add_text_box(slide, 0.5, 0.3, 9, 0.6, "📋 솔루션 포트폴리오", 28, True, COLORS['dark'], PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 0.8, 9, 0.4, "3개 핵심 솔루션 소개", 14, False, COLORS['gray'], PP_ALIGN.CENTER)

    # 솔루션 1
    add_shape_with_text(slide, 0.5, 1.4, 0.6, 0.6, "1", COLORS['primary'], COLORS['white'], 20, True)
    add_text_box(slide, 1.3, 1.4, 8, 0.4, "PDF 문서 모바일 HTML 변환", 16, True, COLORS['dark'])
    add_text_box(slide, 1.3, 1.8, 8, 0.6, "교회 주보, 선거 공보물, 공공기관 소식지, 강의자료, 전공 서적을\n모바일에 최적화된 HTML로 자동 변환", 11, False, COLORS['gray'])
    add_text_box(slide, 1.3, 2.4, 8, 0.3, "AI OCR | 8개국어 | 반응형", 10, False, COLORS['primary'])

    # 솔루션 2
    add_shape_with_text(slide, 0.5, 2.9, 0.6, 0.6, "2", COLORS['orange'], COLORS['white'], 20, True)
    add_text_box(slide, 1.3, 2.9, 8, 0.4, "특수 직업군 외국어 학습기", 16, True, COLORS['dark'])
    add_text_box(slide, 1.3, 3.3, 8, 0.6, "치과, 한방병원, 호텔, 레스토랑, 병원 간호사, 면세점, 미용실 등\n직군별 맞춤 다국어 학습 솔루션", 11, False, COLORS['gray'])
    add_text_box(slide, 1.3, 3.9, 8, 0.3, "7개 직군 | 실무 표현 | 음성 지원", 10, False, COLORS['orange'])

    # 솔루션 3
    add_shape_with_text(slide, 0.5, 4.4, 0.6, 0.6, "3", COLORS['accent'], COLORS['white'], 20, True)
    add_text_box(slide, 1.3, 4.4, 8, 0.4, "통합 유통관리 시스템", 16, True, COLORS['dark'])
    add_text_box(slide, 1.3, 4.8, 8, 0.6, "도매점-소매점-소비자 연계, 온/오프라인 통합\n신개념 유통환경 최적화 시스템", 11, False, COLORS['gray'])
    add_text_box(slide, 1.3, 5.4, 8, 0.3, "B2B2C | 카카오 연동 | 자동 정산", 10, False, COLORS['accent'])

def create_slide_3_pdf_overview(prs):
    """슬라이드 3: PDF 변환 솔루션 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    add_shape_with_text(slide, 0.3, 0.2, 0.7, 0.7, "📱", COLORS['primary'], COLORS['white'], 24)
    add_text_box(slide, 1.1, 0.25, 7, 0.4, "PDF 문서 모바일 HTML 변환", 22, True, COLORS['dark'])
    add_text_box(slide, 1.1, 0.6, 7, 0.3, "AI 기반 자동 레이아웃 분석 & 반응형 변환", 12, False, COLORS['gray'])

    # 핵심 가치
    add_text_box(slide, 0.3, 1.1, 4.2, 0.35, "🎯 핵심 가치", 13, True, COLORS['primary'])
    features1 = "✓ A4 PDF를 모바일 최적화 HTML로 자동 변환\n✓ Claude Vision AI 기반 레이아웃 분석\n✓ 카테고리별 맞춤 템플릿 자동 적용\n✓ 다크모드, 폰트 조절, 공유 기능 내장"
    add_text_box(slide, 0.3, 1.45, 4.2, 1.3, features1, 10, False, COLORS['dark'])

    # 지원 문서 유형
    add_text_box(slide, 5.0, 1.1, 4.5, 0.35, "📊 지원 문서 유형", 13, True, COLORS['primary'])
    features2 = "✓ 교회 주보 - 8개국어 번역, 성경/찬송 팝업\n✓ 선거 공보물 - 후보자 정보, 공약 구조화\n✓ 공공기관 소식지 - 뉴스레터 디지털화\n✓ 강의자료/전공서적 - 학습 콘텐츠 최적화"
    add_text_box(slide, 5.0, 1.45, 4.5, 1.3, features2, 10, False, COLORS['dark'])

    # 변환 프로세스
    add_text_box(slide, 0.3, 2.9, 9, 0.35, "⚙️ 변환 프로세스", 13, True, COLORS['primary'])

    steps = ["📄\nPDF 업로드", "🤖\nAI 분석", "🌍\n다국어 번역", "🎨\n템플릿 적용", "📱\nHTML 생성"]
    for i, step in enumerate(steps):
        add_shape_with_text(slide, 0.4 + i*1.9, 3.3, 1.5, 0.9, step, COLORS['light'], COLORS['dark'], 9)
        if i < 4:
            add_text_box(slide, 1.9 + i*1.9, 3.6, 0.4, 0.3, "→", 16, True, COLORS['primary'], PP_ALIGN.CENTER)

    # 하이라이트 박스
    add_shape_with_text(slide, 0.5, 4.5, 2.8, 0.8, "변환 소요 시간\n3~5초", COLORS['primary'], COLORS['white'], 12, True)
    add_shape_with_text(slide, 3.6, 4.5, 2.8, 0.8, "지원 언어\n8개국어", COLORS['accent'], COLORS['white'], 12, True)
    add_shape_with_text(slide, 6.7, 4.5, 2.8, 0.8, "모바일 최적화\n95%+", COLORS['orange'], COLORS['white'], 12, True)

def create_slide_4_church(prs):
    """슬라이드 4: 교회 주보 상세"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    add_shape_with_text(slide, 0.3, 0.2, 0.7, 0.7, "⛪", COLORS['primary'], COLORS['white'], 24)
    add_text_box(slide, 1.1, 0.25, 7, 0.4, "교회 주보 변환 시스템", 22, True, COLORS['dark'])
    add_text_box(slide, 1.1, 0.6, 7, 0.3, "8개국어 자동 번역 + 성경/찬송 연동", 12, False, COLORS['gray'])

    # 4개 카드
    add_text_box(slide, 0.3, 1.1, 4.4, 0.35, "🌍 다국어 지원", 13, True, COLORS['primary'])
    lang = "✓ 한국어 - 원본 텍스트\n✓ 영어 - English\n✓ 중국어 - 简体中文\n✓ 일본어 - 日本語\n✓ 인도네시아어, 스페인어, 러시아어, 프랑스어"
    add_text_box(slide, 0.3, 1.45, 4.4, 1.4, lang, 10, False, COLORS['dark'])

    add_text_box(slide, 5.0, 1.1, 4.4, 0.35, "📖 인터랙티브 기능", 13, True, COLORS['primary'])
    inter = "✓ 성경 구절 클릭 → 다국어 본문 팝업\n✓ 찬송가 번호 클릭 → 다국어 가사 팝업\n✓ 설교 본문 전체 번역\n✓ 오늘의 양식 다국어 지원"
    add_text_box(slide, 5.0, 1.45, 4.4, 1.2, inter, 10, False, COLORS['dark'])

    add_text_box(slide, 0.3, 3.0, 4.4, 0.35, "📱 주보 섹션 구조", 13, True, COLORS['primary'])
    section = "✓ 오늘의 말씀 (컴팩트 아코디언)\n✓ 예배 안내 (공통순서 + 개별 카드)\n✓ 생명의 말씀 (설교 아코디언)\n✓ 금주의 찬양 / 교회소식 / 헌금 안내"
    add_text_box(slide, 0.3, 3.35, 4.4, 1.2, section, 10, False, COLORS['dark'])

    add_text_box(slide, 5.0, 3.0, 4.4, 0.35, "🛠️ 기술 스택", 13, True, COLORS['primary'])
    tech = "✓ Python Flask + Claude Vision API\n✓ PyMuPDF (PDF 처리)\n✓ 다국어 폴백: 선택어 → 영어 → 한국어\n✓ 반응형 CSS + 다크모드"
    add_text_box(slide, 5.0, 3.35, 4.4, 1.2, tech, 10, False, COLORS['dark'])

def create_slide_5_other_docs(prs):
    """슬라이드 5: 기타 PDF 변환"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    add_shape_with_text(slide, 0.3, 0.2, 0.7, 0.7, "📄", COLORS['primary'], COLORS['white'], 24)
    add_text_box(slide, 1.1, 0.25, 7, 0.4, "다양한 문서 유형 지원", 22, True, COLORS['dark'])
    add_text_box(slide, 1.1, 0.6, 7, 0.3, "선거 공보물 / 공공기관 소식지 / 강의자료", 12, False, COLORS['gray'])

    cards = [
        ("🗳️ 선거 공보물", "✓ 후보자 정보 자동 추출\n✓ 공약 구조화 (카드 형태)\n✓ 정당별 테마 색상 적용\n✓ 비교 보기 기능", 0.3, 1.1),
        ("🏛️ 공공기관 소식지", "✓ 뉴스 카테고리별 분류\n✓ 이미지 갤러리 자동 생성\n✓ 연락처/일정 카드화\n✓ SNS 공유 기능", 5.0, 1.1),
        ("📚 강의자료 / 전공서적", "✓ 챕터별 네비게이션\n✓ 수식/도표 렌더링\n✓ 하이라이트 & 메모 기능\n✓ 목차 자동 생성", 0.3, 3.0),
        ("💡 공통 기능", "✓ 반응형 디자인 (모바일 우선)\n✓ 다크모드 / 폰트 크기 조절\n✓ 오프라인 저장 (PWA)\n✓ URL 공유 & QR 코드", 5.0, 3.0),
    ]

    for title, content, left, top in cards:
        add_text_box(slide, left, top, 4.4, 0.35, title, 13, True, COLORS['primary'])
        add_text_box(slide, left, top + 0.35, 4.4, 1.3, content, 10, False, COLORS['dark'])

def create_slide_6_language_overview(prs):
    """슬라이드 6: 외국어 학습기 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    add_shape_with_text(slide, 0.3, 0.2, 0.7, 0.7, "🌍", COLORS['orange'], COLORS['white'], 24)
    add_text_box(slide, 1.1, 0.25, 7, 0.4, "특수 직업군 외국어 학습기", 22, True, COLORS['dark'])
    add_text_box(slide, 1.1, 0.6, 7, 0.3, "직군별 실무 표현 집중 학습", 12, False, COLORS['gray'])

    add_text_box(slide, 0.3, 1.1, 4.4, 0.35, "🎯 솔루션 특징", 13, True, COLORS['orange'])
    features = "✓ 직업군별 맞춤 커리큘럼\n✓ 실무에서 바로 쓰는 표현 중심\n✓ 다국어 지원 (영/중/일/동남아)\n✓ 음성 발음 가이드\n✓ 상황별 롤플레이 시나리오"
    add_text_box(slide, 0.3, 1.45, 4.4, 1.5, features, 10, False, COLORS['dark'])

    add_text_box(slide, 5.0, 1.1, 4.4, 0.35, "👥 지원 직업군 (7개)", 13, True, COLORS['orange'])
    jobs = "🦷 치과 - 진료/치료 안내 표현\n🌿 한방병원 - 한의학 용어, 침/뜸 설명\n🏨 호텔 - 체크인/컨시어지 서비스\n🍽️ 레스토랑 - 주문/서빙/결제\n🏥 병원 간호사 - 환자 응대/의료 용어\n🛍️ 면세점 - 상품 안내/세금 환급\n💇 미용실 - 스타일 상담/시술 설명"
    add_text_box(slide, 5.0, 1.45, 4.4, 2.0, jobs, 10, False, COLORS['dark'])

    add_text_box(slide, 0.3, 3.6, 9, 0.35, "📱 학습 흐름", 13, True, COLORS['orange'])
    steps = ["👤\n직군 선택", "📚\n상황별 학습", "🔊\n발음 연습", "🎭\n롤플레이", "📊\n학습 현황"]
    for i, step in enumerate(steps):
        add_shape_with_text(slide, 0.4 + i*1.9, 4.0, 1.5, 0.9, step, COLORS['light'], COLORS['dark'], 9)
        if i < 4:
            add_text_box(slide, 1.9 + i*1.9, 4.3, 0.4, 0.3, "→", 16, True, COLORS['orange'], PP_ALIGN.CENTER)

def create_slide_7_language_examples(prs):
    """슬라이드 7: 외국어 학습기 예시"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape_with_text(slide, 0.3, 0.2, 0.7, 0.7, "📖", COLORS['orange'], COLORS['white'], 24)
    add_text_box(slide, 1.1, 0.25, 7, 0.4, "직군별 학습 콘텐츠 예시", 22, True, COLORS['dark'])
    add_text_box(slide, 1.1, 0.6, 7, 0.3, "실무 상황 기반 다국어 표현", 12, False, COLORS['gray'])

    examples = [
        ("🦷 치과 - 진료 안내", '✓ "오늘 어디가 불편하세요?"\n✓ "스케일링을 받으신 적 있으세요?"\n✓ "마취가 풀릴 때까지 식사를 피해주세요"\n✓ "3개월 후 정기 검진을 권장드립니다"', 0.3, 1.1),
        ("🏨 호텔 - 체크인/서비스", '✓ "예약 확인 도와드리겠습니다"\n✓ "조식은 7시부터 10시까지입니다"\n✓ "Wi-Fi 비밀번호를 안내해 드릴게요"\n✓ "체크아웃은 오전 11시입니다"', 5.0, 1.1),
        ("🛍️ 면세점 - 쇼핑 안내", '✓ "여권을 보여주시겠어요?"\n✓ "면세 한도는 $800입니다"\n✓ "택스 리펀드 받으실 수 있습니다"\n✓ "포장해 드릴까요?"', 0.3, 2.9),
        ("🏥 병원 간호사 - 환자 응대", '✓ "알레르기가 있으신가요?"\n✓ "혈압을 측정하겠습니다"\n✓ "주사 맞으실 때 따끔할 수 있어요"\n✓ "처방전은 1층 약국에서 받으세요"', 5.0, 2.9),
    ]

    for title, content, left, top in examples:
        add_text_box(slide, left, top, 4.4, 0.35, title, 12, True, COLORS['orange'])
        add_text_box(slide, left, top + 0.35, 4.4, 1.2, content, 9, False, COLORS['dark'])

    add_shape_with_text(slide, 0.5, 4.6, 2.8, 0.7, "지원 직업군\n7개", COLORS['orange'], COLORS['white'], 11, True)
    add_shape_with_text(slide, 3.6, 4.6, 2.8, 0.7, "학습 표현\n500+", COLORS['primary'], COLORS['white'], 11, True)
    add_shape_with_text(slide, 6.7, 4.6, 2.8, 0.7, "지원 언어\n5개", COLORS['accent'], COLORS['white'], 11, True)

def create_slide_8_demo(prs):
    """슬라이드 8: 데모 링크"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape_with_text(slide, 0.3, 0.2, 0.7, 0.7, "🔗", COLORS['primary'], COLORS['white'], 24)
    add_text_box(slide, 1.1, 0.25, 7, 0.4, "라이브 데모", 22, True, COLORS['dark'])
    add_text_box(slide, 1.1, 0.6, 7, 0.3, "실제 서비스 체험하기", 12, False, COLORS['gray'])

    add_text_box(slide, 0.5, 1.2, 4.0, 0.4, "📱 StudySnap Mobile", 16, True, COLORS['dark'], PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 1.6, 4.0, 0.3, "PDF 변환 + 외국어 학습기 통합 데모", 10, False, COLORS['gray'], PP_ALIGN.CENTER)
    add_shape_with_text(slide, 0.8, 2.0, 3.4, 0.5, "studysnap-mobile.netlify.app", COLORS['primary'], COLORS['white'], 11, True)

    add_text_box(slide, 5.0, 1.2, 4.4, 0.35, "🖥️ 데모 체험 내용", 13, True, COLORS['dark'])
    demo = "✓ 교회 주보 PDF 변환 체험\n✓ 8개국어 언어 전환 테스트\n✓ 성경/찬송 팝업 기능\n✓ 다크모드 & 반응형 확인\n✓ 직업군별 외국어 학습"
    add_text_box(slide, 5.0, 1.55, 4.4, 1.4, demo, 10, False, COLORS['dark'])

    add_text_box(slide, 0.3, 3.0, 9, 0.4, "📋 지원 변환기 목록", 13, True, COLORS['dark'])

    headers = ["변환기", "입력", "출력", "특징"]
    for i, h in enumerate(headers):
        add_shape_with_text(slide, 0.3 + i*2.35, 3.5, 2.3, 0.4, h, COLORS['primary'], COLORS['white'], 10, True)

    rows = [
        ["교회 주보", "PDF (6페이지)", "반응형 HTML", "8개국어, 성경 연동"],
        ["선거 공보물", "PDF", "후보자 카드", "공약 구조화"],
        ["강의자료", "PDF/PPT", "학습 모듈", "챕터 네비게이션"],
        ["외국어 학습", "커리큘럼", "인터랙티브", "7개 직군, 음성"],
    ]
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            color = COLORS['light'] if r % 2 == 0 else COLORS['white']
            add_shape_with_text(slide, 0.3 + c*2.35, 3.9 + r*0.4, 2.3, 0.4, cell, color, COLORS['dark'], 9)

def create_slide_9_kmart_overview(prs):
    """슬라이드 9: K-MART 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape_with_text(slide, 0.3, 0.2, 0.7, 0.7, "🛒", COLORS['accent'], COLORS['white'], 24)
    add_text_box(slide, 1.1, 0.25, 7, 0.4, "통합 유통관리 시스템 (K-MART)", 20, True, COLORS['dark'])
    add_text_box(slide, 1.1, 0.6, 7, 0.3, "도매점-소매점-소비자 연계 B2B2C 플랫폼", 12, False, COLORS['gray'])

    add_text_box(slide, 0.3, 1.1, 4.4, 0.35, "🎯 솔루션 개요", 13, True, COLORS['accent'])
    overview = "✓ 온라인 + 오프라인 통합 유통환경\n✓ 도매점 → 소매점 → 소비자 3단계 연계\n✓ 카카오톡 기반 모바일 주문 시스템\n✓ 실시간 재고/정산 통합 관리"
    add_text_box(slide, 0.3, 1.45, 4.4, 1.2, overview, 10, False, COLORS['dark'])

    add_text_box(slide, 5.0, 1.1, 4.4, 0.35, "👥 사용자 계층", 13, True, COLORS['accent'])
    users = "✓ 도매점 - 상품 등록, 소매점 관리, 통합 정산\n✓ 소매점 - 독립 재고/가격, 자체 마진 설정\n✓ 소비자 - 카카오톡 간편 주문"
    add_text_box(slide, 5.0, 1.45, 4.4, 1.0, users, 10, False, COLORS['dark'])

    add_text_box(slide, 0.3, 2.8, 9, 0.35, "🔄 거래 흐름", 13, True, COLORS['accent'])
    steps = ["🏭\n도매점", "🏪\n소매점", "📲\n카카오 링크", "👤\n소비자", "💰\n자동 정산"]
    for i, step in enumerate(steps):
        add_shape_with_text(slide, 0.4 + i*1.9, 3.2, 1.5, 0.9, step, COLORS['light'], COLORS['dark'], 9)
        if i < 4:
            add_text_box(slide, 1.9 + i*1.9, 3.5, 0.4, 0.3, "→", 16, True, COLORS['accent'], PP_ALIGN.CENTER)

def create_slide_10_kmart_features(prs):
    """슬라이드 10: K-MART 기능 상세"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape_with_text(slide, 0.3, 0.2, 0.7, 0.7, "⚙️", COLORS['accent'], COLORS['white'], 24)
    add_text_box(slide, 1.1, 0.25, 7, 0.4, "K-MART 핵심 기능", 22, True, COLORS['dark'])
    add_text_box(slide, 1.1, 0.6, 7, 0.3, "Django 6.0 기반 통합 관리 시스템", 12, False, COLORS['gray'])

    cards = [
        ("📦 상품 관리", "✓ 무제한 계층 카테고리\n✓ 도매가/소매가 이중 가격\n✓ 소매점 독립 재고\n✓ 엑셀 대량 업로드", 0.3, 1.1),
        ("🛒 주문 관리", "✓ 주문 상태 5단계 관리\n✓ 카카오톡 자동 알림\n✓ 비회원 간편 주문\n✓ 출고 시 재고 자동 차감", 5.0, 1.1),
        ("💰 정산 관리", "✓ 일정산 / 익일정산 선택\n✓ 기간별 매출/매입 자동 집계\n✓ 거래처별 미수금 현황\n✓ 정산 이력 로그", 0.3, 2.8),
        ("📊 대시보드", "✓ 오늘/이번달 매출 통계\n✓ 소매점별 판매 현황\n✓ 재고 부족 알림\n✓ 종합 리포트", 5.0, 2.8),
    ]

    for title, content, left, top in cards:
        add_text_box(slide, left, top, 4.4, 0.35, title, 12, True, COLORS['accent'])
        add_text_box(slide, left, top + 0.35, 4.4, 1.1, content, 10, False, COLORS['dark'])

    add_shape_with_text(slide, 0.5, 4.5, 2.8, 0.7, "거래 유형\n3종", COLORS['accent'], COLORS['white'], 11, True)
    add_shape_with_text(slide, 3.6, 4.5, 2.8, 0.7, "주문 상태\n5단계", COLORS['primary'], COLORS['white'], 11, True)
    add_shape_with_text(slide, 6.7, 4.5, 2.8, 0.7, "정산 방식\n2종", COLORS['orange'], COLORS['white'], 11, True)

def create_slide_11_comparison(prs):
    """슬라이드 11: 솔루션 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape_with_text(slide, 0.3, 0.2, 0.7, 0.7, "📊", COLORS['secondary'], COLORS['white'], 24)
    add_text_box(slide, 1.1, 0.25, 7, 0.4, "솔루션 비교 요약", 22, True, COLORS['dark'])
    add_text_box(slide, 1.1, 0.6, 7, 0.3, "기술 스택 및 특징", 12, False, COLORS['gray'])

    headers = ["구분", "PDF 모바일 변환", "외국어 학습기", "유통관리 시스템"]
    for i, h in enumerate(headers):
        w = 1.8 if i == 0 else 2.4
        left = 0.3 if i == 0 else 2.1 + (i-1)*2.4
        add_shape_with_text(slide, left, 1.1, w, 0.45, h, COLORS['primary'], COLORS['white'], 10, True)

    rows = [
        ["핵심 기능", "PDF → HTML 변환", "직군별 다국어 학습", "B2B2C 통합 관리"],
        ["기술 스택", "Flask, Claude API", "React, TTS API", "Django 6.0, Kakao"],
        ["주요 고객", "교회, 정당, 교육기관", "병원, 호텔, 면세점", "도매/소매 유통업체"],
        ["특징", "8개국어, 성경 연동", "7개 직군, 음성 지원", "3단계 계층, 자동 정산"],
        ["배포 방식", "SaaS / On-premise", "SaaS / 앱", "SaaS / 맞춤 개발"],
        ["현황", "✅ 운영 중", "✅ 운영 중", "✅ 운영 중"],
    ]

    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            w = 1.8 if c == 0 else 2.4
            left = 0.3 if c == 0 else 2.1 + (c-1)*2.4
            color = COLORS['light'] if r % 2 == 0 else COLORS['white']
            font_bold = c == 0
            add_shape_with_text(slide, left, 1.55 + r*0.55, w, 0.55, cell, color, COLORS['dark'], 9, font_bold)

def create_slide_12_contact(prs):
    """슬라이드 12: 연락처"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['dark']
    background.line.fill.background()

    add_text_box(slide, 0.5, 0.8, 9, 0.6, "프로젝트 문의", 32, True, COLORS['white'], PP_ALIGN.CENTER)

    add_text_box(slide, 1.5, 2.0, 2.5, 0.5, "📧", 40, False, COLORS['white'], PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 2.5, 2.5, 0.3, "이메일", 12, False, COLORS['gray'], PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 2.8, 2.5, 0.4, "jmyangkr@gmail.com", 14, True, COLORS['white'], PP_ALIGN.CENTER)

    add_text_box(slide, 3.75, 2.0, 2.5, 0.5, "📞", 40, False, COLORS['white'], PP_ALIGN.CENTER)
    add_text_box(slide, 3.75, 2.5, 2.5, 0.3, "연락처", 12, False, COLORS['gray'], PP_ALIGN.CENTER)
    add_text_box(slide, 3.75, 2.8, 2.5, 0.4, "010-8665-8150", 14, True, COLORS['white'], PP_ALIGN.CENTER)

    add_text_box(slide, 6.0, 2.0, 2.5, 0.5, "🔗", 40, False, COLORS['white'], PP_ALIGN.CENTER)
    add_text_box(slide, 6.0, 2.5, 2.5, 0.3, "데모", 12, False, COLORS['gray'], PP_ALIGN.CENTER)
    add_text_box(slide, 6.0, 2.8, 2.5, 0.4, "studysnap-mobile.netlify.app", 11, True, COLORS['white'], PP_ALIGN.CENTER)

    add_text_box(slide, 0.5, 3.8, 9, 0.4, "맞춤형 솔루션 개발 및 기술 컨설팅 상담 가능",
                 14, False, COLORS['gray'], PP_ALIGN.CENTER)

    add_text_box(slide, 0.5, 4.5, 9, 0.8, "Thank You", 48, True, COLORS['purple'], PP_ALIGN.CENTER)

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    create_slide_1_cover(prs)
    create_slide_2_toc(prs)
    create_slide_3_pdf_overview(prs)
    create_slide_4_church(prs)
    create_slide_5_other_docs(prs)
    create_slide_6_language_overview(prs)
    create_slide_7_language_examples(prs)
    create_slide_8_demo(prs)
    create_slide_9_kmart_overview(prs)
    create_slide_10_kmart_features(prs)
    create_slide_11_comparison(prs)
    create_slide_12_contact(prs)

    output_path = "portfolio_2025.pptx"
    prs.save(output_path)
    print(f"PPTX file created: {output_path}")

if __name__ == "__main__":
    main()
